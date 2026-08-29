"""Aplica o plano de revisão sobre a apresentação do time.

O arquivo de origem é o `_vtime.pptx`, feito pelo time. Este roteiro não
regenera o desenho: ele edita textos e remove slides, para que a identidade
visual azul continue exatamente como está.

O plano vive em `Plano Apresentação MedFlow.md` e pede três coisas: reescrever
alguns slides, fundir tópicos que estão repartidos e reduzir a contagem. As
fusões preservam o assunto — nenhum tema sai do deck, o que sai é a repetição.

Uso:
    python revisar_deck_time.py
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation

BASE = Path(__file__).resolve().parent
ORIGEM = BASE / "EC_Sprint_2_1TSCO_EvidenciasConstrucao_MedFlow_OmegaUrbanTech_vtime.pptx"
DESTINO = BASE / "EC_Sprint_2_1TSCO_EvidenciasConstrucao_MedFlow_OmegaUrbanTech.pptx"

# O número da página fica encostado na margem direita, a 12,13 polegadas.
CANTO_DIREITO = int(11.5 * 914400)
CAPTURAS_DIR = BASE / "capturas"

# Slides absorvidos por outro, na numeração original de 30 slides.
ABSORVIDOS = {
    13,  # Camada Gold -> some no Pipeline, que passa a explicar o pipeline inteiro
    15,  # API pública -> some em Armazenamento e serving
    20,  # Catálogo de indicadores 2 de 2 -> some no 1 de 2
    23,  # Evidências visuais 3 de 3 -> some no 1 de 3
    25,  # Entregáveis técnicos -> some em Validação técnica
    29,  # Evolução -> some em Resultados
}

# O 22 voltou. Ele fora absorvido junto com o 23 quando os três painéis de
# evidência viraram um, mas a evidência da IA não tinha onde aparecer: o slide
# da FlowIA é todo texto, e a captura do APEX — que é entregável — estava
# sobrando na pasta. O 22 já tem a estrutura de dois painéis com legenda, então
# volta reescrito, com a FlowIA respondendo ao vivo e o rastro do APEX.

# (slide na numeração original, texto exato, substituto)
TROCAS: list[tuple[int, str, str]] = []


def troca(slide: int, velho: str, novo: str) -> None:
    TROCAS.append((slide, velho, novo))


# Trocas que valem para **todas** as caixas iguais do slide, não só a primeira.
# Selos repetidos, como os três "VALIDADA" da coluna da direita, precisam disso:
# `troca` para no primeiro frame que casa, de propósito, porque rótulos curtos
# se repetem entre slides e trocar todos por engano estragaria o desenho.
TROCAS_TODAS: list[tuple[int, str, str]] = []


def troca_todas(slide: int, velho: str, novo: str) -> None:
    TROCAS_TODAS.append((slide, velho, novo))


# ---------------------------------------------------------------- slide 4
# O plano pede a linha proativa: as fontes existem antes do problema, mas só
# são lidas depois dele.
troca(4, "Dados públicos existem — leitura integrada, não",
      "Os dados chegam antes do problema. A leitura, depois dele")
troca(
    4,
    "O SUS publica mensalmente um volume massivo de dados administrativos sobre "
    "internações, leitos, diagnósticos e valores. Eles são públicos e abertos, mas "
    "vivem em domínios distintos, com finalidades e temporalidades próprias.",
    "O SUS publica todo mês dados administrativos sobre internações, leitos, "
    "diagnósticos e valores. São públicos, abertos e chegam sozinhos — mas vivem "
    "em domínios distintos, e juntá-los custa tempo. Por isso costumam ser abertos "
    "depois que alguém já reclamou. A pergunta aqui é a inversa: o que esta "
    "competência já mostra.",
)

# ---------------------------------------------------------------- slide 5
troca(
    5,
    "“Quando recebo uma nova competência dos dados públicos hospitalares, quero "
    "localizar rapidamente regiões e hospitais com sinais fora do seu histórico ou "
    "dos seus pares, entender o recorte que concentra o sinal e levar uma hipótese "
    "documentada para investigação com a equipe local.”",
    "“Quando uma competência nova é publicada, quero saber onde olhar antes que o "
    "problema chegue até mim: quais regiões e hospitais saíram do próprio histórico "
    "ou dos pares, qual recorte concentra o sinal, e o que levo como hipótese "
    "verificável para a equipe local — sem precisar montar um projeto de dados "
    "para cada pergunta.”",
)
troca(
    5,
    "Sem uma leitura integrada, cresce o esforço para localizar regiões e hospitais "
    "que merecem investigação — e para distinguir oferta, demanda, deslocamento, "
    "capacidade declarada e resultado administrativo.",
    "Sem leitura integrada, a competência nova só é aberta quando alguém já "
    "reclamou. Aí a pergunta deixa de ser onde investigar e passa a ser como "
    "justificar o que já aconteceu.",
)

# ---------------------------------------------------------------- slide 6
troca(
    6,
    "Precisa priorizar investigações territoriais e hospitalares a partir de dados "
    "administrativos públicos.",
    "Recebe uma competência nova todo mês, com as bases na mão. E costuma "
    "descobrir o problema pelo telefone, não pelo dado.",
)

# ---------------------------------------------------------------- slide 7
troca(
    7,
    "O MedFlow não cria dados novos. Ele conforma fontes públicas com finalidades "
    "diferentes em uma leitura mensal comparável, calcula os indicadores uma única "
    "vez na camada Gold e serve esse resultado por uma API somente leitura — com "
    "fórmula, amostra e limitação sempre ao lado do número.",
    "Antes: cinco bases com finalidades diferentes, cada pergunta virando um "
    "projeto de dados, e a resposta chegando depois do problema. Depois: a "
    "competência publicada vira, no mesmo dia, uma lista de onde olhar primeiro. "
    "O MedFlow não cria dados novos — conforma as fontes uma vez, calcula os "
    "indicadores uma vez na Gold e serve por uma API somente leitura, com fórmula, "
    "amostra e limite ao lado de cada número.",
)

# ---------------------------------------------------------------- slide 8
# "Só revisar título de IA": a dimensão é a de IA generativa, e o rótulo não
# dizia isso. Os números também avançaram desde a montagem do deck.
troca(8, "Linguagem natural", "IA generativa")
troca(8, "Select AI previsto para a Sprint 2",
      "Select AI previsto, sem profile nem validação")
troca(8, "Profile MEDFLOW_GENAI com bateria de 5 perguntas controladas",
      "Profile MEDFLOW_GENAI, bateria de 5 perguntas e assistente FlowIA no produto")
troca(8, "10 endpoints ORDS GET no módulo público api/v1",
      "10 endpoints GET e o POST do assistente no módulo público api/v1")
troca(8, "WebApp React 19 + TypeScript + Vite, público no GitHub Pages",
      "WebApp React 19 + TypeScript + Vite, duas páginas e quatro visões, no GitHub Pages")
troca(8, "5 índices + fluxos, ICSAP, taxas populacionais e CMI real por IPCA",
      "6 índices com o IPE, mais fluxos, ICSAP, taxas populacionais e CMI real por IPCA")

# ---------------------------------------------------------------- slide 9
troca(9, "Quatro visões navegáveis",
      "Duas páginas navegáveis")
troca(9, "Regional, Fluxos, Hospital e Metodologia, com filtros por URL.",
      "Análise contínua com território e hospital, mais a Metodologia; recorte na URL.")
troca(9, "10 endpoints GET publicados por ORDS no módulo api/v1.",
      "10 endpoints GET e o POST do assistente, publicados por ORDS em api/v1.")
troca(9, "Cinco perguntas controladas em linguagem natural",
      "Assistente FlowIA com IA generativa governada")
troca(9, "Select AI sobre o schema MEDFLOW, com SQL de referência.",
      "Select AI sobre o schema MEDFLOW, com SQL auditável e memória de conversa.")
troca(9, "Indicadores calculados e persistidos na Gold",
      "Seis indicadores calculados e persistidos na Gold")

# ---------------------------------------------------------------- slide 10
# O slide 12 já tinha sido corrigido para 12 tabelas e 5 dimensões, mas a
# arquitetura e o fechamento ficaram com o número antigo. Como é o mesmo dado
# dito em três lugares, a divergência aparece para quem folheia o deck.
troca(10, "Indicadores,\n2 dimensões · 7 marts",
      "Indicadores,\n5 dimensões · 7 marts")
troca(10, "Lakehouse, schema\nMEDFLOW, 9 tabelas",
      "Lakehouse, schema\nMEDFLOW, 12 tabelas")
troca(
    10,
    "Views de projeção pura, sem recálculo. Módulo público api/v1 com 10 endpoints "
    "GET, somente leitura.",
    "Views de projeção pura, sem recálculo. Módulo público api/v1 com 10 endpoints "
    "GET e o POST do assistente.",
)
# "Quatro visões" e "duas páginas" convivem no deck sem que ele diga que são
# as duas coisas ao mesmo tempo: duas rotas, quatro visões dentro delas.
troca(
    10,
    "React 19 + TypeScript, publicado no GitHub Pages. Quatro visões consomem "
    "valores prontos da API.",
    "React 19 + TypeScript, publicado no GitHub Pages. Duas páginas e quatro "
    "visões consomem valores prontos da API.",
)

# ---------------------------------------------------------------- slide 11
# O plano manda excluir o bloco de validação pendente.
troca(11, "Validação ainda pendente", "Governança da IA")
troca(
    11,
    "Comparação final das duas perguntas territoriais novas do Select AI, pelo "
    "mesmo processo aplicado às três perguntas originais.",
    "Todo SQL gerado passa por guarda de somente leitura, auditoria em tabela "
    "própria e cota diária. A narrativa só é liberada depois da conferência do SQL.",
)
troca(
    11,
    "O profile e as cinco perguntas existem e respondem; o que falta é o "
    "fechamento homogêneo da evidência comparativa.",
    "A resposta narrativa nunca é a fonte da verdade: ela é rejeitada quando o SQL "
    "diverge ou quando o texto ultrapassa a metodologia.",
)
# A coluna é estreita (3,14 pol) e cabem ~45 caracteres por linha. Passar disso
# quebra o item em duas linhas e o marcador redondo, que tem posição fixa, fica
# pendurado no meio do bloco em vez de alinhado à primeira linha.
troca(11, "WebApp público com quatro visões",
      "WebApp público: duas páginas e quatro visões")
troca(11, "Select AI com profile e bateria de cinco perguntas",
      "Select AI com profile, bateria e FlowIA")

# ---------------------------------------------------------------- slide 12
# Vira o slide único do pipeline: absorve o que a Camada Gold detalhava.
troca(12, "Bronze, Silver e Gold — o que cada camada garante",
      "Bronze, Silver e Gold — o caminho inteiro do dado")
troca(
    12,
    "Fórmulas, denominadores, amostras, benchmarks, fluxos, ICSAP, correção por "
    "IPCA e agregações geográficas — a partir apenas de estruturas contratadas da "
    "Silver.",
    "Sete marts, um grão de análise cada: hospital, hospital/especialidade, "
    "hospital/CID, região mensal e por período, fluxos e ICSAP. Fórmulas, "
    "denominadores, amostras, benchmarks, correção por IPCA e agregações "
    "geográficas, a partir apenas de estruturas contratadas da Silver.",
)
troca(12, "2 + 7", "12")
troca(12, "dimensões geográficas e marts", "tabelas no Oracle: 5 dimensões e 7 marts")
troca(12, "597.725", "597.930")
troca(
    12,
    "É a fonte semântica única dos valores usados pelo banco, pela API e pelas telas.",
    "Fonte semântica única do banco, da API e das telas: 189 colunas comentadas, e "
    "indicador nenhum é recalculado depois daqui.",
)

# ---------------------------------------------------------------- slide 14
# Absorve a API pública: banco, governança e o que ele expõe num slide só.
troca(14, "ARMAZENAMENTO E SERVING", "ARMAZENAMENTO E API")
troca(14, "Oracle Autonomous AI Database 26ai Lakehouse",
      "Oracle 26ai Lakehouse e os onze endpoints públicos")
troca(
    14,
    "Nenhuma tabela Gold publicada diretamente por AutoREST",
    "Nenhuma tabela por AutoREST: api/v1 expõe 10 GET e 1 POST",
)
troca(14, "597.725", "597.930")
troca(14, "linhas carregadas no recorte vigente",
      "linhas carregadas nas 12 tabelas do recorte vigente")
troca(14, "36 / 36", "47 / 47")
troca(14, "métricas do gate Oracle com estado ok",
      "portões do gate Oracle com estado ok")
troca(
    14,
    "O plano Always Free pode hibernar por inatividade. O último estado registrado "
    "em 16/08/2026 confirmou conexão mTLS, integridade da Gold e os dez endpoints "
    "públicos respondendo.",
    "O plano Always Free pode hibernar por inatividade. A conferência de "
    "apresentação roda em um comando e checa conexão, contrato, competência, os "
    "endpoints e o link público — 12 de 12 na última execução.",
)

# ---------------------------------------------------------------- slide 16
# Refeito: o slide falava do Select AI como capacidade solta. Quem a banca vê
# no produto é a FlowIA, e o Select AI é o motor por baixo dela. A coluna da
# esquerda passa a ser o caminho de uma pergunta; a da direita deixa de listar
# as cinco perguntas (que seguem como evidência no slide de validação) e passa
# a dizer o que a avaliação encontrou — inclusive o que ainda não fechou.
troca(16, "Select AI — linguagem natural governada",
      "FlowIA — o assistente que consulta o Select AI")
troca(16, "Papel no projeto", "O que a FlowIA faz")
troca(
    16,
    "Camada controlada de acesso ao modelo Gold e recurso de explicabilidade. "
    "Profile MEDFLOW_GENAI sobre OCI Generative AI por Resource Principal, ampliado "
    "para nove objetos do schema de aplicação.",
    "O assistente do produto. O gestor pergunta em português na tela em que já "
    "está, e a FlowIA manda junto o contexto silencioso — competência, região, "
    "hospital e a análise ativa.",
)
troca(
    16,
    "Não é chat público, não alimenta o WebApp e não substitui Gold, contratos ou "
    "SQL validado.",
    "Pergunta sobre regra de produto a FlowIA responde sozinha, sem custo e sem "
    "modelo. Ao Select AI vai só pergunta de dado — e a resposta nunca substitui "
    "Gold, contrato ou SQL validado.",
)

# As quatro etapas deixam de ser o roteiro de validação do avaliador e passam a
# ser o trajeto que a pergunta do gestor faz até virar resposta.
troca(16, "Processo técnico de validação", "O caminho de uma pergunta")
troca(16, "SQL de referência", "Contexto")
troca(16, "Consulta convencional executada e preservada.",
      "A tela, a competência e as duas rodadas anteriores acompanham a pergunta.")
troca(16, "Inspeção do SQL efetivamente gerado pelo Select AI.",
      "O Select AI gera o SQL sobre o profile MEDFLOW_GENAI, e ele é inspecionado "
      "antes de rodar.")
troca(16, "Comparação", "Guarda")
troca(16, "Filtros, agregações, recortes, ordenação e limite conferidos.",
      "Só executa o que começa por SELECT ou WITH; verbo de escrita derruba a "
      "rodada.")
troca(16, "Narrativa liberada apenas após a coerência do SQL.",
      "A narrativa sai depois do SQL conferido, e passa por varredura de "
      "vocabulário.")
troca(
    16,
    "O resultado é rejeitado quando o SQL diverge ou a narrativa ultrapassa a "
    "metodologia. A resposta narrativa nunca é a fonte da verdade.",
    "Cada pergunta grava uma linha auditável: o SQL, a narrativa e o que a "
    "varredura marcou. Cota de 50 por dia, e a narrativa nunca é a fonte da "
    "verdade.",
)

# A direita passa a ser a avaliação da FlowIA. As três primeiras linhas são
# verdes no desenho e as duas últimas, âmbar — então as corrigidas ficam em
# cima e o que segue aberto, embaixo, sem mexer na cor de nenhuma.
troca(16, "AS CINCO PERGUNTAS REGISTRADAS", "O ROTEIRO FOI TESTADO: 20 DE 20")
troca(16, "Pressão regional", "Bateria de 20 perguntas humanas")
troca(16, "As cinco regiões de saúde com maior IPH médio em 2026.",
      "Curtas, vagas e coloquiais, com o contexto vindo só da tela. Reprovava 0 de "
      "20 em 26/08.")
troca(16, "Mortalidade por especialidade", "Cada resposta conferida na Gold")
troca(
    16,
    "As dez especialidades com maior TMH média, com amostra suficiente e mínimo de "
    "100 linhas hospital-mês.",
    "Um SQL de referência roda na mesma Gold, e os rótulos e a ordem têm de bater "
    "com os do modelo.",
)
troca(16, "IPR por diagnóstico", "Perguntas que devem ser recusadas")
troca(
    16,
    "Os dez diagnósticos com maior IPR médio, com amostra suficiente e ao menos 10 "
    "combinações.",
    "Mais cheio hoje, pior hospital, IPR ruim: a FlowIA responde que o dado não "
    "mede isso.",
)
troca(16, "Evasão intrastadual observada", "Rodada única, ainda por fazer")
troca(
    16,
    "As regiões com maior percentual médio de evasão intrastadual observada em 2026.",
    "Os 20 fecharam em duas rodadas no mesmo dia: a cota diária de 50 perguntas "
    "acabou antes de uma só.",
)
troca(16, "Grupos ICSAP", "Duas perguntas territoriais")
troca(16, "Os dez grupos ICSAP com mais internações de residentes em 2026.",
      "Evasão intrastadual e grupos ICSAP aguardam o fechamento comparativo das "
      "outras três.")
troca(
    16,
    "As três perguntas originais têm evidência completa de SQL de referência, "
    "showsql e narrate comparados. As duas perguntas territoriais adicionadas após "
    "a ampliação do profile aguardam o mesmo fechamento comparativo.",
    "Sete defeitos corrigidos — caixa alta, ordem invertida, mês truncado, buffer, "
    "amostra, ICSAP e competência solta. As correções vivem no prompt do pacote "
    "PL/SQL: valem para a bateria e para o site.",
)
# As três linhas verdes passaram a dizer PASSA. As duas âmbar continuam
# "A REVALIDAR", que é o que a rodada única e as territoriais ainda são.
troca_todas(16, "VALIDADA", "PASSA")

# ---------------------------------------------------------------- slide 17
# O produto deixou de ser quatro visões: a análise virou contínua, com duas
# etapas ancoradas, e a etapa de fluxos saiu por decisão de produto.
troca(17, "WebApp — quatro visões, uma jornada",
      "WebApp — duas páginas, uma investigação contínua")
troca(
    17,
    "React 19 + TypeScript + Vite, publicado no GitHub Pages. As visões consomem "
    "valores prontos da API e aplicam apenas formatação com Intl — sem fórmulas, "
    "faixas ou cortes próprios.",
    "React 19 + TypeScript + Vite, publicado no GitHub Pages. São duas páginas e "
    "quatro visões: Território, Pares e Hospital dividem a análise, ancorados na "
    "mesma leitura, e a Metodologia é a segunda página. As telas consomem valores "
    "prontos da API e aplicam apenas formatação com Intl — sem fórmulas, faixas ou "
    "cortes próprios.",
)
troca(17, "Regional", "Território")
troca(
    17,
    "Competência, macrorregião e região; mapa de IPH por percentis, métricas, "
    "série e ranking com amostra.",
    "Mapa das 62 regiões colorido pelo placar de sinais ou pelo IPH, totais do "
    "recorte com MoM e YoY, e a série mensal da região escolhida.",
)
troca(
    17,
    "Localiza variação territorial e temporal sem converter sinal em conclusão de "
    "desempenho.",
    "Responde onde olhar primeiro, e a série diz se o sinal é do mês ou vem de "
    "antes — sem converter sinal em conclusão de desempenho.",
)
# A visão de fluxos saiu do produto; o cartão passa a ser o da comparação com pares.
troca(17, "Fluxos", "Pares")
troca(17, "A população é atendida no próprio território — e para onde se desloca?",
      "Este número é alto para um hospital deste porte?")
troca(
    17,
    "Residência e destino, atendimento próprio, evasão intrastadual observada, "
    "atração, taxa residente, matriz origem–destino e ICSAP.",
    "Faixa interquartil, mediana e posição do hospital entre pares do mesmo porte, "
    "na mesma região, com os pares nomeados e o critério escrito.",
)
troca(
    17,
    "Separa demanda residente de produção e contextualiza dependência e referência "
    "regional.",
    "O produto sabia dizer o número e não sabia dizer se ele era normal. O porte "
    "entra no critério porque é ele que torna os números comparáveis.",
)
troca(17, "Em quais hospitais, especialidades ou diagnósticos o sinal se concentra?",
      "Em qual hospital, especialidade ou diagnóstico o sinal se concentra?")
troca(
    17,
    "Hospitais da região, série mensal, especialidades e CIDs elegíveis para IPR.",
    "Lista ordenável da região, série mensal, especialidades com o IPE e CIDs "
    "elegíveis para o IPR.",
)
troca(
    17,
    "Aprofunda volume, tendência, mortalidade, valor e permanência com recorte e "
    "amostra.",
    "Aprofunda volume, tendência, mortalidade, valor e permanência ante os pares, "
    "sempre com recorte e amostra declarados.",
)
troca(
    17,
    "Cobertura, competência mais recente, fórmulas, cortes, flags, estados nulos, "
    "reconciliações, fontes e limitações.",
    "Reconciliação e limites em primeiro plano; cobertura, fórmulas, cortes, "
    "estados de ausência, definições e fontes logo abaixo.",
)

# ---------------------------------------------------------------- slides fundidos
# O contador de partes some junto com as partes.
troca(19, "CATÁLOGO DE INDICADORES  ·  1 DE 2", "CATÁLOGO DE INDICADORES")
troca(21, "EVIDÊNCIAS VISUAIS  ·  1 DE 3", "EVIDÊNCIAS VISUAIS")
troca(21, "Visão Regional e visão Fluxos", "O produto em execução")
troca(21, "Regional", "Território")
troca(
    21,
    "Mapa de IPH por percentis, série mensal e ranking de regiões com amostra "
    "declarada. O gestor localiza onde há um sinal fora do padrão antes de escolher "
    "onde aprofundar. Os cortes de percentil e a amostra aparecem ao lado do mapa "
    "para impedir leitura de desempenho.",
    "Mapa das 62 regiões colorido pelo placar de sinais: o tom mais claro é o melhor "
    "do recorte, o mais escuro é o pior, e a escala se estica até o máximo observado. "
    "O gestor vê de longe onde há sinal fora do padrão, e a legenda diz que é placar "
    "para priorizar investigação, não nota de qualidade.",
)
troca(
    21,
    "Matriz origem–destino entre as 62 regiões de saúde, com atendimento próprio, "
    "evasão intrastadual observada, atração assistencial e composição ICSAP. É aqui "
    "que demanda residente e produção hospitalar deixam de ser confundidas.",
    "Comparação com pares do mesmo porte, na mesma região: faixa interquartil, "
    "mediana e a posição do hospital, com o critério escrito e os pares nomeados. "
    "É aqui que o número deixa de ser um valor solto e passa a ter uma régua — e a "
    "régua fixa o porte, porque é ele que torna os números comparáveis.",
)
troca(21, "Fluxos", "Hospital")

# ---------------------------------------------------------------- slide 22
# Reescrito como a evidência da IA: à esquerda a FlowIA respondendo no site
# publicado, à direita o rastro auditado no APEX.
troca(22, "EVIDÊNCIAS VISUAIS  ·  2 DE 3", "EVIDÊNCIA DA IA GENERATIVA")
troca(22, "Visão Hospital e visão Metodologia",
      "A FlowIA respondendo, e o rastro auditado no APEX")
troca(22, "Hospital", "FlowIA no site publicado")
troca(22, "Metodologia", "Select AI no APEX")
troca(
    22,
    "Lista dos hospitais da região, série mensal por unidade, quebra por "
    "especialidade e CIDs elegíveis para IPR. Cada valor traz a amostra e a "
    "elegibilidade — combinações que não atingem os cortes mínimos aparecem como "
    "insuficientes, não como zero.",
    "Pergunta feita ao vivo no site, em 29/08: “quem segura o paciente por mais "
    "tempo?”. O texto não nomeia indicador, tabela nem corte — a FlowIA leu a "
    "tela, escolheu permanência média, aplicou o corte de amostra e fixou a "
    "competência. O SQL fica a um clique.",
)
troca(
    22,
    "Cobertura, competência mais recente, fórmulas, cortes, flags, estados nulos, "
    "reconciliações, fontes e limitações. É a tela que torna o produto auditável: "
    "qualquer número exibido pode ser rastreado até a fórmula e a fonte.",
    "A mesma camada pelo APEX, onde cada pergunta vira uma linha auditada: o SQL "
    "que o modelo gerou, a narrativa e o que a varredura de vocabulário marcou. É "
    "o rastro que sustenta a resposta — a narrativa nunca é a fonte da verdade.",
)
troca(
    22,
    "Os testes de frontend confrontam com Playwright os valores renderizados no "
    "DOM contra os dados da camada Gold, incluindo os estados de erro, ausência e "
    "contingência.",
    "As duas telas passam pelo mesmo pacote PL/SQL — guarda de somente leitura, "
    "corte de amostra, cota diária e auditoria em tabela própria. A bateria de 20 "
    "perguntas humanas testa esse mesmo caminho.",
)

# ---------------------------------------------------------------- slide 24
troca(24, "36 / 36 ok", "47 / 47 ok")
troca(24, "3 de 5 fechadas", "3 de 5 fechadas")

# ---------------------------------------------------------------- slide 30
# O placar de fechamento repetia os números defasados da arquitetura. São os
# mesmos três dados do slide 10 e do 12, e é o último slide que a banca vê.
troca(30, "9", "12")
# 11 e não 10: são 10 GET mais o POST do assistente, como o slide 13 já dizia
# no título. O rótulo fica só "endpoints" — a caixa é estreita demais para a
# conta, e ela já está escrita no slide que apresenta a API.
troca(30, "10", "11")


def texto_dos_slides(prs) -> dict[int, list]:
    """Todos os frames de texto por slide, incluindo os de dentro de tabelas."""
    mapa: dict[int, list] = {}
    for i, slide in enumerate(prs.slides, 1):
        frames = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                frames.append(shape.text_frame)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        frames.append(cell.text_frame)
        mapa[i] = frames
    return mapa


def substituir(frame, velho: str, novo: str) -> bool:
    """Troca o texto preservando a formatação da primeira corrida.

    Um parágrafo do PowerPoint costuma vir picado em várias `runs`, cada uma com
    a própria formatação. Escrever na primeira e esvaziar as demais mantém fonte,
    corpo e cor do texto original, que é o que preserva a identidade do deck.
    """
    if frame.text.strip() != velho.strip():
        return False
    paragrafos = frame.paragraphs
    primeiro = paragrafos[0]
    if not primeiro.runs:
        return False
    primeiro.runs[0].text = novo
    for run in primeiro.runs[1:]:
        run.text = ""
    for paragrafo in paragrafos[1:]:
        for run in paragrafo.runs:
            run.text = ""
    return True


def remover_slide(prs, indice_zero: int) -> None:
    """Remove um slide pelo índice, junto com a relação que o referencia."""
    slides = prs.slides._sldIdLst  # noqa: SLF001 - a API pública não expõe remoção
    ids = list(slides)
    alvo = ids[indice_zero]
    prs.part.drop_rel(alvo.rId)
    slides.remove(alvo)


# Painel de evidência -> captura que entra nele, na numeração original.
CAPTURAS = {
    "Print da visão Regional": "jornada-1-territorio.png",
    "Print da visão Fluxos": "jornada-4-pares.png",
    # Slide 22, a evidência da IA. A da FlowIA é o recorte da conversa real
    # capturada por capturar_flowia.mjs; a do APEX é a de 23/08, que já existia
    # e não tinha slide desde que os painéis de evidência foram fundidos.
    "Print da visão Hospital": "flowia-ao-vivo-conversa.png",
    "Print da visão Metodologia": "apex_select_ai_real.png",
}


def encaixar(caixa: tuple[int, int, int, int], imagem: Path) -> tuple[int, int, int, int]:
    """Maior retângulo com a proporção da imagem que cabe dentro da caixa.

    Encaixa sem cortar: uma captura de tela cortada perde justamente a borda
    onde costuma estar o rótulo que prova o que a tela mostra.
    """
    from PIL import Image

    esquerda, topo, largura, altura = caixa
    with Image.open(imagem) as img:
        proporcao = img.width / img.height
    if largura / altura > proporcao:
        nova_altura = altura
        nova_largura = int(altura * proporcao)
    else:
        nova_largura = largura
        nova_altura = int(largura / proporcao)
    return (
        esquerda + (largura - nova_largura) // 2,
        topo + (altura - nova_altura) // 2,
        nova_largura,
        nova_altura,
    )


def inserir_capturas(prs) -> list[str]:
    """Põe cada captura no painel que a anuncia e apaga o texto de espera."""
    relatorio = []
    for slide in prs.slides:
        marcadores = [
            shape
            for shape in slide.shapes
            if shape.has_text_frame
            and any(chave in shape.text_frame.text for chave in CAPTURAS)
        ]
        for marcador in marcadores:
            chave = next(k for k in CAPTURAS if k in marcador.text_frame.text)
            imagem = CAPTURAS_DIR / CAPTURAS[chave]
            if not imagem.exists():
                relatorio.append(f"AUSENTE {imagem.name}")
                continue
            # O painel é o retângulo que envolve o marcador: mesma coluna, um
            # pouco mais à esquerda e mais alto.
            painel = min(
                (
                    shape
                    for shape in slide.shapes
                    if shape.left is not None
                    and shape.left < marcador.left
                    and shape.top < marcador.top
                    and shape.width > marcador.width
                ),
                key=lambda shape: (marcador.left - shape.left) + (marcador.top - shape.top),
            )
            caixa = (painel.left, painel.top, painel.width, painel.height)
            esquerda, topo, largura, altura = encaixar(caixa, imagem)
            slide.shapes.add_picture(str(imagem), esquerda, topo, largura, altura)
            for paragrafo in marcador.text_frame.paragraphs:
                for run in paragrafo.runs:
                    run.text = ""
            relatorio.append(
                f"{imagem.name} em {esquerda / 914400:.2f},{topo / 914400:.2f} "
                f"({largura / 914400:.2f}x{altura / 914400:.2f})"
            )
    return relatorio


POL = 914400


def _perto(valor: int, alvo_pol: float, folga_pol: float = 0.03) -> bool:
    return abs(valor - alvo_pol * POL) <= folga_pol * POL


def ajustes_de_layout(prs) -> list[str]:
    """Acerta espaçamentos que a revisão de texto desencaixou.

    Roda depois da remoção dos slides, então trabalha na numeração final. Move
    caixas e altura de linha; não mexe em cor, fonte nem no desenho dos cartões.
    """
    feitos = []

    # -------------------------------------------------------------- slide 4
    # A linha de chamada ficava à mesma distância do parágrafo e dos cartões,
    # sem dizer a qual dos dois pertence. Aproximá-la dos cartões, que é o que
    # ela introduz, resolve pela proximidade em vez de por mais espaço.
    for shape in prs.slides[3].shapes:
        if shape.has_text_frame and shape.text_frame.text.startswith("Para responder"):
            shape.top = int(2.50 * POL)
            feitos.append("slide 4: linha de chamada aproximada dos cartões")

    # ------------------------------------------------------------- slide 15
    # Os quatro cartões foram desenhados para um texto maior do que o que
    # sobrou depois da revisão, e abriram um vão no meio. Sobe o bloco de
    # baixo e encurta o cartão na mesma medida.
    subida = int(0.45 * POL)
    movidas = 0
    for shape in prs.slides[14].shapes:
        if not _perto(shape.left, 0.60) and not _perto(shape.left, 3.71) \
           and not _perto(shape.left, 6.82) and not _perto(shape.left, 9.93) \
           and not _perto(shape.left, 0.82) and not _perto(shape.left, 3.93) \
           and not _perto(shape.left, 7.04) and not _perto(shape.left, 10.15):
            continue
        if _perto(shape.top, 4.72) or _perto(shape.top, 4.84) or _perto(shape.top, 5.06):
            shape.top -= subida
            movidas += 1
        elif _perto(shape.top, 2.04) and _perto(shape.height, 4.24):
            shape.height = int(3.62 * POL)
    if movidas:
        feitos.append(f"slide 15: {movidas} caixas subiram e os 4 cartões encurtaram")

    # ------------------------------------------------------------- slide 17
    # A tabela ocupava só a metade de cima, com meio slide vazio embaixo. As
    # linhas crescem para preencher o espaço que já era dela.
    for shape in prs.slides[16].shapes:
        if getattr(shape, "has_table", False):
            linhas = shape.table.rows
            linhas[0].height = int(0.55 * POL)
            for linha in list(linhas)[1:]:
                linha.height = int(0.70 * POL)
            feitos.append("slide 17: tabela esticada para preencher o slide")

    return feitos


def main() -> int:
    prs = Presentation(ORIGEM)
    frames = texto_dos_slides(prs)

    aplicadas, ausentes = 0, []
    for slide, velho, novo in TROCAS:
        if velho == novo:
            continue
        if any(substituir(frame, velho, novo) for frame in frames.get(slide, [])):
            aplicadas += 1
        else:
            ausentes.append((slide, velho[:60]))

    for slide, velho, novo in TROCAS_TODAS:
        if velho == novo:
            continue
        # `sum` em vez de `any`: consome tudo, em vez de parar no primeiro acerto.
        if sum(substituir(frame, velho, novo) for frame in frames.get(slide, [])):
            aplicadas += 1
        else:
            ausentes.append((slide, velho[:60]))

    # O número de rodapé é um texto solto no canto inferior direito, não um campo
    # automático: sem renumerar, o slide 13 continuaria anunciando "14".
    restantes = [n for n in range(1, len(prs.slides.__iter__.__self__._sldIdLst) + 1)]  # noqa: SLF001
    restantes = [n for n in range(1, 31) if n not in ABSORVIDOS]
    for nova, original in enumerate(restantes, 1):
        for frame in frames[original]:
            if frame.text.strip() == str(original) and frame._parent.left > CANTO_DIREITO:  # noqa: SLF001
                substituir(frame, str(original), str(nova))

    for indice in sorted(ABSORVIDOS, reverse=True):
        remover_slide(prs, indice - 1)

    inseridas = inserir_capturas(prs)
    ajustados = ajustes_de_layout(prs)

    prs.save(DESTINO)

    pedidas = len([t for t in TROCAS + TROCAS_TODAS if t[1] != t[2]])
    print(f"trocas aplicadas: {aplicadas} de {pedidas}")
    for slide, trecho in ausentes:
        print(f"  NAO ENCONTRADO no slide {slide}: {trecho}")
    print(f"slides: {len(Presentation(ORIGEM).slides)} -> {len(prs.slides)}")
    for linha in inseridas:
        print(f"  captura: {linha}")
    for linha in ajustados:
        print(f"  layout: {linha}")
    print(f"gravado em {DESTINO.name}")
    return 1 if ausentes else 0


if __name__ == "__main__":
    raise SystemExit(main())
