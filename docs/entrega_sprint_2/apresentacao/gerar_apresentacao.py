"""Gera a apresentação final da Sprint 2 do MedFlow.

O arquivo de origem do time está em `99_arquivo/02_oracle_medflow/apresentacao/`.
Esta versão usa capturas reais do produto e o arco aprovado de 21 slides.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent
CAPTURAS = BASE / "capturas"
SAIDA = BASE / "EC_Sprint_2_1TSCO_EvidenciasConstrucao_MedFlow_OmegaUrbanTech.pptx"

PRODUTO_URL = "https://lucas-d-s-1.github.io/medflow/"
REPOSITORIO_URL = "https://github.com/Lucas-D-S-1/medflow"
# O vídeo é produzido pelo time. Preencher somente com a URL pública definitiva.
VIDEO_URL = ""

W = 13.333
H = 7.5

PAPER = "F5F5EF"
WHITE = "FFFFFF"
INK = "123F36"
TEXT = "38564F"
MUTED = "6E817B"
GREEN = "08765E"
MINT = "DCEFE7"
MINT_DARK = "72CBB1"
LINE = "CBD9D2"
AMBER = "D88922"
AMBER_SOFT = "F8E6C7"
RED_SOFT = "F4DAD5"
NAVY = "172A4A"

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"
FONT_MONO = "Aptos Mono"


def rgb(hexadecimal: str) -> RGBColor:
    return RGBColor.from_string(hexadecimal)


def fundo(slide, cor: str = PAPER) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(cor)


def caixa(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str | None = None,
    radius: bool = False,
):
    tipo = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(tipo, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(0.8 if line else 0)
    return shape


def linha(slide, x1: float, y1: float, x2: float, y2: float, cor: str = LINE, largura: float = 1.5):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = rgb(cor)
    shape.line.width = Pt(largura)
    return shape


def area_link(slide, x: float, y: float, w: float, h: float, destino: str):
    """Sobrepõe uma área clicável transparente sem tingir o texto de azul."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.fill.background()
    shape.click_action.hyperlink.address = destino
    return shape


def texto(
    slide,
    conteudo: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    tamanho: float = 18,
    cor: str = TEXT,
    negrito: bool = False,
    fonte: str = FONT_BODY,
    alinhar=PP_ALIGN.LEFT,
    ancora=MSO_ANCHOR.TOP,
    margem: float = 0,
    hyperlink: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margem)
    tf.vertical_anchor = ancora
    p = tf.paragraphs[0]
    p.alignment = alinhar
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = conteudo
    run.font.name = fonte
    run.font.size = Pt(tamanho)
    run.font.bold = negrito
    if hyperlink:
        run.hyperlink.address = hyperlink
        run.font.underline = False
    run.font.color.rgb = rgb(cor)
    return shape


def lista(
    slide,
    itens: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    tamanho: float = 16,
    cor: str = TEXT,
    marcador: str = "—",
    espacamento: float = 10,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{marcador}  {item}"
        p.font.name = FONT_BODY
        p.font.size = Pt(tamanho)
        p.font.color.rgb = rgb(cor)
        p.space_after = Pt(espacamento)
        p.line_spacing = 1.05
    return shape


def rotulo(slide, conteudo: str, x: float, y: float, w: float = 4.0, cor: str = GREEN):
    return texto(
        slide,
        conteudo.upper(),
        x,
        y,
        w,
        0.28,
        tamanho=9.5,
        cor=cor,
        negrito=True,
    )


def titulo(slide, conteudo: str, numero: int, *, escuro: bool = False, subtitulo: str | None = None):
    cor = WHITE if escuro else INK
    rotulo(slide, f"MEDFLOW · {numero:02d}", 0.65, 0.38, cor=MINT_DARK if escuro else GREEN)
    texto(slide, conteudo, 0.65, 0.78, 12.0, 1.05, tamanho=30, cor=cor, negrito=True, fonte=FONT_HEAD)
    if subtitulo:
        texto(slide, subtitulo, 0.68, 1.75, 11.7, 0.52, tamanho=14, cor=MINT if escuro else MUTED)


def rodape(slide, numero: int, *, escuro: bool = False):
    cor = MINT_DARK if escuro else MUTED
    texto(slide, "FIAP × ORACLE · ÔMEGA URBAN TECH · 2026", 0.68, 7.12, 6.5, 0.2, tamanho=7.5, cor=cor)
    texto(slide, str(numero), 12.35, 7.08, 0.32, 0.22, tamanho=8, cor=cor, alinhar=PP_ALIGN.RIGHT)


def imagem_capa(
    slide,
    caminho: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    vertical: str = "center",
    horizontal: str = "center",
    borda: str | None = LINE,
):
    if borda:
        caixa(slide, x - 0.025, y - 0.025, w + 0.05, h + 0.05, fill=WHITE, line=borda, radius=True)
    with Image.open(caminho) as im:
        iw, ih = im.size
    proporcao_imagem = iw / ih
    proporcao_alvo = w / h
    pic = slide.shapes.add_picture(str(caminho), Inches(x), Inches(y), Inches(w), Inches(h))
    if proporcao_imagem > proporcao_alvo:
        total = 1 - proporcao_alvo / proporcao_imagem
        if horizontal == "left":
            pic.crop_left, pic.crop_right = 0, total
        elif horizontal == "right":
            pic.crop_left, pic.crop_right = total, 0
        else:
            pic.crop_left = pic.crop_right = total / 2
    elif proporcao_imagem < proporcao_alvo:
        total = 1 - proporcao_imagem / proporcao_alvo
        if vertical == "top":
            pic.crop_top, pic.crop_bottom = 0, total
        elif vertical == "bottom":
            pic.crop_top, pic.crop_bottom = total, 0
        else:
            pic.crop_top = pic.crop_bottom = total / 2
    return pic


def passo(slide, numero: str, nome: str, pergunta: str, x: float, y: float, w: float, *, ativo: bool = False):
    cor = GREEN if ativo else INK
    caixa(slide, x, y, w, 1.42, fill=WHITE, line=cor if ativo else LINE, radius=True)
    texto(slide, numero, x + 0.2, y + 0.16, 0.5, 0.34, tamanho=12, cor=cor, negrito=True)
    texto(slide, nome, x + 0.2, y + 0.48, w - 0.4, 0.35, tamanho=18, cor=INK, negrito=True, fonte=FONT_HEAD)
    texto(slide, pergunta, x + 0.2, y + 0.9, w - 0.4, 0.4, tamanho=10.5, cor=MUTED)


def metrica(slide, valor: str, legenda: str, x: float, y: float, w: float, *, cor: str = GREEN):
    texto(slide, valor, x, y, w, 0.5, tamanho=28, cor=cor, negrito=True, fonte=FONT_HEAD)
    texto(slide, legenda, x, y + 0.5, w, 0.45, tamanho=10.5, cor=MUTED)


def criar() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # 01 — capa
    s = prs.slides.add_slide(blank)
    fundo(s, INK)
    caixa(s, 7.3, 0, 6.03, 7.5, fill=PAPER)
    imagem_capa(s, CAPTURAS / "produto-abertura.png", 7.42, 0.65, 5.72, 6.2, horizontal="left", borda=None)
    rotulo(s, "Enterprise Challenge 2026 · FIAP × Oracle", 0.68, 0.55, 6.0, cor=MINT_DARK)
    texto(s, "Uma nova\ncompetência chegou.", 0.68, 1.18, 6.15, 1.65, tamanho=36, cor=WHITE, negrito=True, fonte=FONT_HEAD)
    texto(s, "Onde investigar primeiro?", 0.7, 3.0, 5.9, 0.72, tamanho=25, cor=MINT_DARK, negrito=True, fonte=FONT_HEAD)
    texto(s, "Do sinal territorial à hipótese: com fonte, amostra e limite à vista.", 0.7, 3.93, 5.9, 0.82, tamanho=16, cor=MINT)
    caixa(s, 0.7, 4.92, 5.85, 0.62, fill=GREEN, line=MINT_DARK, radius=True)
    texto(s, "Med = saúde · Flow = fluxo do dado + fluxo assistencial", 0.92, 5.1, 5.42, 0.25, tamanho=10.5, cor=WHITE, negrito=True, alinhar=PP_ALIGN.CENTER)
    texto(s, "Sprint 2 · MVP público · 30 competências · São Paulo", 0.7, 5.82, 5.9, 0.38, tamanho=11, cor=MINT)
    texto(s, "Ômega Urban Tech", 0.7, 6.55, 3.0, 0.3, tamanho=10, cor=WHITE, negrito=True)

    # 02 — cenário de uso
    s = prs.slides.add_slide(blank)
    fundo(s)
    titulo(s, "Quem prioriza. E quem ganha quando ele acerta", 2, subtitulo="Cenário de uso hipotético: não é evidência de entrevista ou teste com usuário.")
    texto(s, "06/2026", 0.75, 2.5, 2.5, 0.85, tamanho=40, cor=GREEN, negrito=True, fonte=FONT_HEAD)
    texto(s, "Uma nova competência dos dados hospitalares foi publicada.", 0.78, 3.38, 4.0, 0.95, tamanho=22, cor=INK, negrito=True, fonte=FONT_HEAD)
    texto(s, "Persona: gestor ou analista regional do SUS.", 0.78, 4.62, 4.0, 0.42, tamanho=12.5, cor=GREEN, negrito=True)
    texto(s, "Beneficiário final: a população atendida.", 0.78, 5.12, 4.0, 0.42, tamanho=12.5, cor=INK, negrito=True)
    linha(s, 5.0, 2.45, 5.0, 6.25, GREEN, 2)
    for y, n, chamada, detalhe in [
        (2.52, "01", "Onde há um sinal?", "62 regiões não cabem numa leitura manual coerente."),
        (3.72, "02", "O que dá contexto?", "Produção, residência, fluxo e capacidade não são sinônimos."),
        (4.92, "03", "O que levar para a conversa?", "Uma hipótese verificável, com amostra, fórmula e limite."),
    ]:
        caixa(s, 5.35, y, 0.52, 0.52, fill=INK, radius=True)
        texto(s, n, 5.35, y + 0.08, 0.52, 0.25, tamanho=9, cor=WHITE, negrito=True, alinhar=PP_ALIGN.CENTER)
        texto(s, chamada, 6.05, y - 0.03, 4.9, 0.38, tamanho=18, cor=INK, negrito=True, fonte=FONT_HEAD)
        texto(s, detalhe, 6.05, y + 0.42, 5.6, 0.52, tamanho=13, cor=MUTED)
    texto(s, "A decisão continua humana. O produto organiza a investigação para quem precisa priorizar.", 5.35, 6.12, 6.55, 0.52, tamanho=14, cor=GREEN, negrito=True)
    rodape(s, 2)

    # 03 — problema
    s = prs.slides.add_slide(blank)
    fundo(s, WHITE)
    titulo(s, "A pergunta é simples. O caminho até ela, não.", 3)
    caixa(s, 0.8, 2.23, 11.75, 0.86, fill=INK, radius=True)
    texto(s, "“Onde devo investigar primeiro?”", 1.1, 2.42, 11.1, 0.42, tamanho=24, cor=WHITE, negrito=True, fonte=FONT_HEAD, alinhar=PP_ALIGN.CENTER)
    fontes = [
        ("SIH/SUS", "RELACIONAL · internações"),
        ("CNES atual", "JSON · leitos declarados"),
        ("DATASUS/MS", "CSV · território + ICSAP"),
        ("IBGE", "CSV / EXTERNAL TABLE · população + IPCA"),
    ]
    for i, (nome, detalhe) in enumerate(fontes):
        x = 0.9 + i * 3.0
        texto(s, nome, x, 3.55, 2.55, 0.35, tamanho=17, cor=INK, negrito=True, fonte=FONT_HEAD)
        texto(s, detalhe, x, 3.98, 2.55, 0.52, tamanho=9.5, cor=MUTED)
        if i < 3:
            texto(s, "+", x + 2.55, 3.64, 0.3, 0.3, tamanho=18, cor=GREEN, negrito=True, alinhar=PP_ALIGN.CENTER)
    linha(s, 1.0, 4.58, 12.2, 4.58, LINE, 1)
    conceitos = ["demanda residente", "produção", "fluxo", "capacidade declarada", "resultado administrativo"]
    for i, c in enumerate(conceitos):
        x = 0.82 + i * 2.47
        caixa(s, x, 5.05, 2.2, 0.7, fill=MINT if i in (0, 2) else PAPER, line=LINE, radius=True)
        texto(s, c, x + 0.12, 5.24, 1.96, 0.28, tamanho=11, cor=INK, negrito=True, alinhar=PP_ALIGN.CENTER)
    texto(s, "Se esses conceitos se misturam, o painel responde rápido — e errado.", 0.85, 6.28, 11.6, 0.42, tamanho=18, cor=GREEN, negrito=True, fonte=FONT_HEAD, alinhar=PP_ALIGN.CENTER)
    rodape(s, 3)

    # 04 — jornada
    s = prs.slides.add_slide(blank)
    fundo(s)
    titulo(s, "O MedFlow organiza a investigação, não a decisão", 4, subtitulo="A jornada é a espinha dorsal da apresentação e do produto.")
    xs = [0.75, 3.9, 7.05, 10.2]
    dados = [
        ("01", "Localizar", "Qual região merece atenção?"),
        ("02", "Contextualizar", "O fluxo ajuda a explicar?"),
        ("03", "Aprofundar", "Onde o sinal se concentra?"),
        ("04", "Formular", "Qual hipótese levar adiante?"),
    ]
    for i, (n, nome, pergunta) in enumerate(dados):
        passo(s, n, nome, pergunta, xs[i], 2.55, 2.55, ativo=i == 0)
        if i < 3:
            linha(s, xs[i] + 2.55, 3.26, xs[i + 1], 3.26, GREEN, 2)
    texto(s, "REGIONAL", 1.25, 4.55, 1.55, 0.28, tamanho=9, cor=GREEN, negrito=True, alinhar=PP_ALIGN.CENTER)
    texto(s, "FLUXOS", 4.45, 4.55, 1.45, 0.28, tamanho=9, cor=GREEN, negrito=True, alinhar=PP_ALIGN.CENTER)
    texto(s, "HOSPITAL + CID", 7.38, 4.55, 1.9, 0.28, tamanho=9, cor=GREEN, negrito=True, alinhar=PP_ALIGN.CENTER)
    texto(s, "CONVERSA LOCAL", 10.48, 4.55, 1.95, 0.28, tamanho=9, cor=GREEN, negrito=True, alinhar=PP_ALIGN.CENTER)
    caixa(s, 1.18, 5.35, 10.9, 0.72, fill=WHITE, line=LINE, radius=True)
    texto(s, "Sinal → contexto → recorte → hipótese documentada", 1.5, 5.56, 10.25, 0.28, tamanho=19, cor=INK, negrito=True, fonte=FONT_HEAD, alinhar=PP_ALIGN.CENTER)
    texto(s, "O MedFlow não prescreve ação clínica nem prova causalidade.", 2.0, 6.35, 9.3, 0.34, tamanho=12, cor=MUTED, alinhar=PP_ALIGN.CENTER)
    rodape(s, 4)

    # 05 — produto
    s = prs.slides.add_slide(blank)
    fundo(s, INK)
    rotulo(s, "O produto", 0.65, 0.45, cor=MINT_DARK)
    texto(s, "A interface começa\npela pergunta do gestor.", 0.65, 1.0, 3.95, 1.45, tamanho=25, cor=WHITE, negrito=True, fonte=FONT_HEAD)
    texto(s, "Sem vocabulário de pipeline na entrada. Fonte, competência e limites aparecem no contexto da própria decisão.", 0.68, 2.5, 3.85, 1.65, tamanho=13, cor=MINT)
    texto(s, "ABRIR PRODUTO PÚBLICO ↗", 0.68, 4.35, 3.4, 0.32, tamanho=11, cor=MINT_DARK, negrito=True)
    area_link(s, 0.68, 4.3, 3.4, 0.42, PRODUTO_URL)
    imagem_capa(s, CAPTURAS / "produto-abertura.png", 4.78, 0.55, 8.0, 6.45, horizontal="left", borda=MINT_DARK)
    rodape(s, 5, escuro=True)

    # 06 — regional
    s = prs.slides.add_slide(blank)
    fundo(s)
    titulo(s, "1 / Localizar um sinal sem transformá-lo em diagnóstico", 6)
    imagem_capa(s, CAPTURAS / "jornada-1-regional.png", 0.7, 2.15, 8.0, 3.86, borda=LINE)
    rotulo(s, "Exemplo · Jundiaí · 06/2026", 9.15, 2.18, 3.2)
    metrica(s, "67,2%", "IPH estimado — não é ocupação real", 9.15, 2.72, 3.25)
    metrica(s, "4.613", "internações novas", 9.15, 3.82, 3.25)
    metrica(s, "13", "hospitais com produção", 9.15, 4.92, 3.25)
    texto(s, "Mapa + amostra + denominador ficam na mesma leitura.", 9.15, 5.98, 3.25, 0.72, tamanho=11.5, cor=INK, negrito=True)
    rodape(s, 6)

    # 07 — fluxos
    s = prs.slides.add_slide(blank)
    fundo(s, WHITE)
    titulo(s, "2 / Contextualizar o território antes de julgar a oferta", 7)
    rotulo(s, "Jundiaí · residência e destino", 0.72, 2.22, 3.4)
    metrica(s, "94,4%", "atendidos no próprio território", 0.72, 2.72, 3.55)
    metrica(s, "5,6%", "evasão intrastadual observada", 0.72, 3.86, 3.55, cor=AMBER)
    metrica(s, "14,1%", "atração assistencial", 0.72, 5.0, 3.55)
    texto(s, "O RD-SP não observa saídas para outras UFs. O limite acompanha o número.", 0.72, 6.05, 3.45, 0.78, tamanho=9.5, cor=MUTED)
    imagem_capa(s, CAPTURAS / "jornada-2-fluxos.png", 4.55, 2.1, 8.15, 4.65, vertical="center", borda=LINE)
    rodape(s, 7)

    # 08 — hospital e CID
    s = prs.slides.add_slide(blank)
    fundo(s)
    titulo(s, "3 / Aprofundar o hospital dentro daquela região", 8, subtitulo="O HU Hospital Universitário está dentro de Jundiaí, o mesmo caso acompanhado desde o slide 6.")
    imagem_capa(s, CAPTURAS / "jornada-3-hospitais.png", 0.68, 2.14, 5.98, 4.25, vertical="top", borda=LINE)
    imagem_capa(s, CAPTURAS / "jornada-4-diagnosticos.png", 6.88, 2.14, 5.78, 4.25, vertical="top", borda=LINE)
    caixa(s, 3.3, 6.25, 6.75, 0.56, fill=INK, radius=True)
    texto(s, "Hospital → especialidade → CID → benchmark sem o próprio hospital", 3.48, 6.41, 6.4, 0.25, tamanho=11.5, cor=WHITE, negrito=True, alinhar=PP_ALIGN.CENTER)
    rodape(s, 8)

    # 09 — hipótese
    s = prs.slides.add_slide(blank)
    fundo(s, INK)
    titulo(s, "4 / Formular a hipótese e tornar a próxima pergunta rastreável", 9, escuro=True)
    rotulo(s, "Hipótese de investigação · período agregado", 0.72, 2.28, 5.0, cor=MINT_DARK)
    texto(s, "No HU Hospital Universitário,\nbronquiolite por vírus sincicial respiratório apresenta IPR 1,71.", 0.72, 2.8, 7.0, 1.52, tamanho=20, cor=WHITE, negrito=True, fonte=FONT_HEAD)
    caixa(s, 8.2, 2.33, 4.35, 3.1, fill=WHITE, line=MINT_DARK, radius=True)
    metrica(s, "5,38 dias", "permanência média do hospital · 201 internações", 8.55, 2.72, 3.7)
    metrica(s, "3,14 dias", "benchmark · 83 internações em 6 hospitais", 8.55, 3.9, 3.7)
    texto(s, "Próxima pergunta", 0.72, 4.72, 2.4, 0.3, tamanho=10, cor=MINT_DARK, negrito=True)
    texto(s, "O sinal se repete por competência e encontra explicação no perfil assistencial local?", 0.72, 5.14, 6.8, 1.08, tamanho=18, cor=MINT, negrito=True, fonte=FONT_HEAD)
    texto(s, "IPR não mede qualidade e não prova ineficiência.", 8.55, 5.72, 3.7, 0.42, tamanho=11, cor=MINT)
    rodape(s, 9, escuro=True)

    # 10 — confiança
    s = prs.slides.add_slide(blank)
    fundo(s, WHITE)
    titulo(s, "Confiança também é dizer o que o dado não sustenta", 10)
    imagem_capa(s, CAPTURAS / "metodologia-evidencia.png", 6.15, 2.12, 6.55, 4.62, horizontal="left", borda=LINE)
    pares = [
        ("IPH", "pressão estimada", "ocupação real"),
        ("TMH", "mortalidade administrativa", "qualidade ou causalidade"),
        ("CMI", "valor aprovado", "custo contábil"),
        ("Fluxo", "movimento observado em SP", "evasão total do residente"),
    ]
    for i, (ind, correto, incorreto) in enumerate(pares):
        y = 2.3 + i * 1.0
        texto(s, ind, 0.75, y, 0.65, 0.3, tamanho=12, cor=GREEN, negrito=True)
        texto(s, correto, 1.55, y - 0.04, 1.9, 0.48, tamanho=10.5, cor=INK, negrito=True)
        texto(s, f"≠ {incorreto}", 3.55, y - 0.02, 2.15, 0.48, tamanho=8.5, cor=AMBER)
        linha(s, 0.75, y + 0.55, 5.7, y + 0.55, LINE, 0.8)
    texto(s, "Fórmula, amostra e limite são parte da interface — não rodapé defensivo.", 0.78, 6.25, 4.95, 0.7, tamanho=12, cor=INK, negrito=True)
    rodape(s, 10)

    # 11 — arquitetura
    s = prs.slides.add_slide(blank)
    fundo(s)
    titulo(s, "Arquitetura: o que roda, o que demonstra e o que é horizonte", 11, subtitulo="O WebApp é determinístico; o Select AI é uma demonstração controlada sobre a mesma Gold.")
    etapas = [
        ("FONTES", "SIH · CNES · MS · IBGE", PAPER),
        ("BRONZE", "preserva + SHA-256", WHITE),
        ("SILVER", "conforma conceitos", WHITE),
        ("GOLD", "calcula uma vez", MINT),
        ("ORACLE 26ai", "9 tabelas", WHITE),
        ("ORDS", "10 GET", WHITE),
        ("WEBAPP", "4 visões", MINT),
    ]
    for i, (nome, desc, fill) in enumerate(etapas):
        x = 0.45 + i * 1.82
        caixa(s, x, 2.55, 1.5, 1.25, fill=fill, line=GREEN if nome == "GOLD" else LINE, radius=True)
        texto(s, nome, x + 0.08, 2.82, 1.34, 0.25, tamanho=10, cor=GREEN if nome in ("GOLD", "WEBAPP") else INK, negrito=True, alinhar=PP_ALIGN.CENTER)
        texto(s, desc, x + 0.08, 3.22, 1.34, 0.3, tamanho=8.5, cor=MUTED, alinhar=PP_ALIGN.CENTER)
        if i < len(etapas) - 1:
            linha(s, x + 1.5, 3.17, x + 1.82, 3.17, GREEN, 1.6)
    linha(s, 7.85, 3.83, 7.85, 4.45, GREEN, 1.7)
    caixa(s, 6.62, 4.45, 2.48, 1.15, fill=INK, line=INK, radius=True)
    texto(s, "SELECT AI", 6.78, 4.7, 2.16, 0.28, tamanho=12, cor=MINT_DARK, negrito=True, alinhar=PP_ALIGN.CENTER)
    texto(s, "pergunta → SQL → guarda", 6.78, 5.1, 2.16, 0.27, tamanho=9, cor=WHITE, alinhar=PP_ALIGN.CENTER)
    caixa(s, 9.45, 4.45, 2.7, 1.15, fill=AMBER_SOFT, line=AMBER, radius=True)
    texto(s, "APEX", 9.62, 4.7, 2.35, 0.28, tamanho=12, cor=INK, negrito=True, alinhar=PP_ALIGN.CENTER)
    texto(s, "demo visível + auditável", 9.62, 5.1, 2.35, 0.27, tamanho=9, cor=TEXT, alinhar=PP_ALIGN.CENTER)
    linha(s, 9.1, 5.03, 9.45, 5.03, AMBER, 1.6)
    caixa(s, 0.78, 6.18, 3.6, 0.55, fill=MINT, line=GREEN, radius=True)
    texto(s, "RODA · Gold + ORDS + WebApp", 0.94, 6.34, 3.28, 0.22, tamanho=9.5, cor=INK, negrito=True, alinhar=PP_ALIGN.CENTER)
    caixa(s, 4.68, 6.18, 3.9, 0.55, fill=AMBER_SOFT, line=AMBER, radius=True)
    texto(s, "DEMONSTRA · Select AI + APEX", 4.84, 6.34, 3.58, 0.22, tamanho=9.5, cor=INK, negrito=True, alinhar=PP_ALIGN.CENTER)
    caixa(s, 8.88, 6.18, 3.6, 0.55, fill=WHITE, line=LINE, radius=True)
    texto(s, "HORIZONTE · validação com gestores", 9.04, 6.34, 3.28, 0.22, tamanho=9.5, cor=MUTED, negrito=True, alinhar=PP_ALIGN.CENTER)
    rodape(s, 11)

    # 12 — técnicas
    s = prs.slides.add_slide(blank)
    fundo(s, WHITE)
    titulo(s, "Técnicas: padrões, agrupamentos e a decisão sobre previsão", 12)
    tecnicas = [
        ("Modelo dimensional", "grão explícito por mart evita médias e joins ambíguos"),
        ("Comparação sazonal", "2026 contra os mesmos meses de 2024 e 2025"),
        ("Benchmark elegível", "hospital comparado a pares regionais, excluindo a si próprio"),
        ("Padrões e agrupamentos", "matriz origem–destino, 19 grupos ICSAP e distribuição do IPR"),
        ("Correção por IPCA", "CMI nominal e real preservados; valor aprovado não é custo"),
    ]
    for i, (nome, desc) in enumerate(tecnicas):
        y = 2.02 + i * 0.76
        texto(s, f"0{i + 1}", 0.82, y, 0.55, 0.32, tamanho=10, cor=GREEN, negrito=True)
        texto(s, nome, 1.55, y - 0.06, 3.05, 0.35, tamanho=15.5, cor=INK, negrito=True, fonte=FONT_HEAD)
        texto(s, desc, 4.75, y - 0.02, 7.1, 0.42, tamanho=11.5, cor=TEXT)
        linha(s, 1.52, y + 0.43, 12.0, y + 0.43, LINE, 0.8)
    caixa(s, 0.82, 5.92, 11.18, 0.86, fill=INK, radius=True)
    texto(s, "Não usamos ML preditivo no MVP. A série cobre só 30 competências e recebe revisões. Primeiro validamos definições e uso; previsão exige janela e alvo acordados.", 1.08, 6.12, 10.66, 0.48, tamanho=10.5, cor=WHITE, negrito=True, alinhar=PP_ALIGN.CENTER)
    rodape(s, 12)

    # 13 — Oracle
    s = prs.slides.add_slide(blank)
    fundo(s)
    titulo(s, "Oracle vivo: e o que acontece quando ele não está", 13)
    metrica(s, "597.725", "linhas Gold carregadas", 0.78, 2.1, 2.65)
    metrica(s, "9", "tabelas no schema MEDFLOW", 3.55, 2.1, 2.55)
    metrica(s, "10", "endpoints GET públicos", 6.25, 2.1, 2.55)
    metrica(s, "26.1.3", "APEX disponível", 9.05, 2.1, 2.45)
    caixa(s, 0.78, 3.45, 5.72, 2.55, fill=INK, radius=True)
    rotulo(s, "Oracle no ar · entrega principal", 1.08, 3.78, 3.6, cor=MINT_DARK)
    texto(s, "GET /ords/medflow/api/v1/status", 1.08, 4.2, 4.65, 0.3, tamanho=12, cor=WHITE, negrito=True, fonte=FONT_MONO)
    area_link(s, 1.08, 4.15, 4.65, 0.42, "https://gf68e03b2a30d55-medflow.adb.sa-saopaulo-1.oraclecloudapps.com/ords/medflow/api/v1/status")
    texto(s, 'status: ok · source: oracle-live\ndata_through: 2026-06 · contract: 0.3.0', 1.08, 4.72, 4.92, 0.64, tamanho=10.5, cor=MINT, fonte=FONT_MONO)
    texto(s, "GitHub Pages → ORDS, sem servidor intermediário.", 1.08, 5.55, 4.92, 0.28, tamanho=10.5, cor=WHITE, negrito=True)
    caixa(s, 6.82, 3.45, 5.72, 2.55, fill=WHITE, line=AMBER, radius=True)
    rotulo(s, "Oracle indisponível · contingência", 7.12, 3.78, 3.9, cor=AMBER)
    texto(s, "10 snapshots rotulados", 7.12, 4.2, 4.8, 0.38, tamanho=20, cor=INK, negrito=True, fonte=FONT_HEAD)
    lista(s, ["um por endpoint público", "mesma versão do contrato", "modo identificado; sem misturar live e snapshot"], 7.12, 4.78, 4.9, 1.05, tamanho=10.5, cor=TEXT, marcador="✓", espacamento=4)
    texto(s, "A fonte exibida faz parte da evidência.", 2.5, 6.38, 8.3, 0.4, tamanho=15, cor=GREEN, negrito=True, alinhar=PP_ALIGN.CENTER)
    rodape(s, 13)

    # 14 — validação
    s = prs.slides.add_slide(blank)
    fundo(s, INK)
    titulo(s, "Validação: cada número fecha com a Gold", 14, escuro=True)
    texto(s, "8.403.103", 0.78, 2.25, 4.15, 0.78, tamanho=43, cor=MINT_DARK, negrito=True, fonte=FONT_HEAD)
    texto(s, "comparações campo a campo", 0.82, 3.1, 3.5, 0.35, tamanho=14, cor=WHITE, negrito=True)
    texto(s, "0 divergências no método executado", 0.82, 3.58, 3.85, 0.35, tamanho=13, cor=MINT)
    linha(s, 5.05, 2.25, 5.05, 5.85, MINT_DARK, 1)
    provas = [
        ("185", "testes Python"),
        ("30 + 2", "Playwright herméticos + live"),
        ("42", "checagens OpenAPI"),
        ("31.792", "comparações no módulo público"),
    ]
    for i, (valor, legenda) in enumerate(provas):
        x = 5.55 + (i % 2) * 3.25
        y = 2.28 + (i // 2) * 1.62
        texto(s, valor, x, y, 2.85, 0.55, tamanho=28, cor=WHITE, negrito=True, fonte=FONT_HEAD)
        texto(s, legenda, x, y + 0.62, 2.85, 0.36, tamanho=11, cor=MINT)
    caixa(s, 5.55, 5.63, 6.65, 0.62, fill=MINT, radius=True)
    texto(s, "Contrato → API → DOM → Gold", 5.75, 5.81, 6.25, 0.26, tamanho=13, cor=INK, negrito=True, alinhar=PP_ALIGN.CENTER)
    texto(s, "Consistência técnica não prova utilidade, impacto ou causalidade.", 0.82, 6.43, 11.4, 0.38, tamanho=13, cor=WHITE, negrito=True, alinhar=PP_ALIGN.CENTER)
    rodape(s, 14, escuro=True)

    # 15 — Select AI de apresentação
    s = prs.slides.add_slide(blank)
    fundo(s, WHITE)
    titulo(s, "Select AI de apresentação: uma pergunta que podemos conferir", 15)
    caixa(s, 0.75, 2.18, 3.25, 3.85, fill=INK, radius=True)
    rotulo(s, "Pergunta", 1.02, 2.48, 1.4, cor=MINT_DARK)
    texto(s, "Quais as cinco regiões\nde saúde com maior IPH\nmédio em 2026?", 1.02, 2.93, 2.72, 2.08, tamanho=14, cor=WHITE, negrito=True, fonte=FONT_HEAD)
    texto(s, "linguagem natural", 1.02, 5.4, 2.2, 0.3, tamanho=10, cor=MINT)
    caixa(s, 4.28, 2.18, 4.25, 3.85, fill=PAPER, line=LINE, radius=True)
    rotulo(s, "SQL gerado", 4.58, 2.48, 1.8)
    sql = 'SELECT região, AVG(IPH)\nFROM Gold\nWHERE ano = 2026\nGROUP BY região\nORDER BY AVG(IPH) DESC\nFETCH FIRST 5 ROWS ONLY'
    texto(s, sql, 4.58, 2.95, 3.6, 2.22, tamanho=13, cor=INK, fonte=FONT_MONO)
    texto(s, "guarda somente leitura", 4.58, 5.4, 2.6, 0.3, tamanho=10, cor=GREEN, negrito=True)
    caixa(s, 8.82, 2.18, 3.8, 3.85, fill=MINT, line=GREEN, radius=True)
    rotulo(s, "Resultado conferido", 9.12, 2.48, 2.4)
    lista(s, ["Limeira", "Franco da Rocha", "Jundiaí", "São José do Rio Preto", "Alto Vale do Paraíba"], 9.12, 2.95, 3.1, 2.3, tamanho=13, cor=INK, marcador="✓", espacamento=6)
    texto(s, "mesmos 5 rótulos, mesma ordem", 9.12, 5.4, 3.0, 0.34, tamanho=10, cor=GREEN, negrito=True)
    caixa(s, 3.1, 6.35, 7.15, 0.47, fill=INK, radius=True)
    texto(s, "6 de 8 perguntas com referência coincidiram exatamente", 3.25, 6.47, 6.85, 0.22, tamanho=10.5, cor=WHITE, negrito=True, alinhar=PP_ALIGN.CENTER)
    rodape(s, 15)

    # 16 — Select AI técnico: limites e APEX
    s = prs.slides.add_slide(blank)
    fundo(s, WHITE)
    titulo(s, "Select AI técnico: execução, limites medidos e prova real", 16)
    imagem_capa(s, CAPTURAS / "apex_select_ai_real.png", 0.68, 2.12, 5.55, 4.28, vertical="top", borda=LINE)
    rotulo(s, "Captura real · MEDFLOW_DEMO · App 100 · Página 1", 0.7, 6.55, 5.55, cor=MUTED)
    caixa(s, 6.55, 2.05, 5.88, 0.58, fill=INK, radius=True)
    texto(s, "showsql → comparação → narrate", 6.78, 2.22, 5.42, 0.25, tamanho=12, cor=WHITE, negrito=True, fonte=FONT_MONO, alinhar=PP_ALIGN.CENTER)
    problemas = [
        ("Ranking mensal", "agrega tarde", "SQL de referência"),
        ("Vocabulário", "‘ocupação’ ≠ IPH", "aviso metodológico"),
        ("Conversação", "perde contexto", "pergunta autossuficiente"),
    ]
    for i, (nome, achado, resposta) in enumerate(problemas):
        y = 2.9 + i * 0.85
        caixa(s, 6.58, y, 0.42, 0.42, fill=AMBER, radius=True)
        texto(s, str(i + 1), 6.58, y + 0.08, 0.42, 0.2, tamanho=8.5, cor=WHITE, negrito=True, alinhar=PP_ALIGN.CENTER)
        texto(s, nome, 7.2, y - 0.02, 1.62, 0.3, tamanho=12, cor=INK, negrito=True, fonte=FONT_HEAD)
        texto(s, achado, 8.88, y, 1.45, 0.3, tamanho=9.5, cor=AMBER, negrito=True)
        texto(s, resposta, 10.42, y, 1.82, 0.34, tamanho=9.5, cor=GREEN, negrito=True)
        linha(s, 7.2, y + 0.52, 12.25, y + 0.52, LINE, 0.7)
    caixa(s, 6.58, 5.52, 5.82, 0.9, fill=MINT, line=GREEN, radius=True)
    texto(s, "2 respostas válidas no rastro · export versionado", 6.84, 5.72, 5.3, 0.28, tamanho=11.5, cor=INK, negrito=True)
    texto(s, "db/apex/05_aplicacao_medflow_select_ai.sql", 6.84, 6.08, 5.05, 0.22, tamanho=8.7, cor=TEXT, fonte=FONT_MONO)
    texto(s, "ABRIR GUIA APEX ↗", 9.2, 6.62, 3.15, 0.24, tamanho=9.5, cor=GREEN, negrito=True, alinhar=PP_ALIGN.RIGHT)
    area_link(s, 9.2, 6.55, 3.15, 0.38, f"{REPOSITORIO_URL}/tree/main/db/apex")
    rodape(s, 16)

    # 17 — resultados e fronteira
    s = prs.slides.add_slide(blank)
    fundo(s)
    titulo(s, "Resultado e fronteira: entregamos investigação, não impacto", 17)
    rotulo(s, "O que existe", 0.8, 2.2, 2.2)
    metrica(s, "30", "competências mensais", 0.8, 2.72, 2.2)
    metrica(s, "62", "regiões de saúde", 3.0, 2.72, 2.2)
    metrica(s, "655", "hospitais", 5.2, 2.72, 2.2)
    metrica(s, "4", "visões públicas", 7.4, 2.72, 2.2)
    metrica(s, "13", "perguntas Select AI", 9.6, 2.72, 2.4)
    linha(s, 0.8, 4.05, 12.1, 4.05, LINE, 1)
    rotulo(s, "O que ainda não foi provado", 0.8, 4.48, 3.5, cor=AMBER)
    lista(s, ["compreensão por gestores reais", "redução de tempo ou retrabalho", "mudança de decisão ou impacto assistencial", "causalidade clínica"], 0.8, 4.95, 5.2, 1.55, tamanho=13, cor=TEXT, marcador="×", espacamento=6)
    caixa(s, 6.55, 4.5, 5.55, 1.55, fill=WHITE, line=GREEN, radius=True)
    texto(s, "Resultado defensável", 6.88, 4.82, 4.9, 0.35, tamanho=18, cor=INK, negrito=True, fonte=FONT_HEAD)
    texto(s, "Um MVP público, rastreável e tecnicamente reconciliado para priorizar onde aprofundar.", 6.88, 5.28, 4.78, 0.78, tamanho=11.5, cor=TEXT)
    rodape(s, 17)

    # 18 — evolução e feedback
    s = prs.slides.add_slide(blank)
    fundo(s, WHITE)
    titulo(s, "Da Sprint 1 até aqui: feedback transformado em evidência", 18)
    caixa(s, 0.82, 2.12, 5.55, 1.45, fill=AMBER_SOFT, line=AMBER, radius=True)
    rotulo(s, "Feedback da Sprint 1", 1.12, 2.43, 2.5, cor=AMBER)
    texto(s, "Aprofundar os diferenciais", 1.12, 2.9, 4.8, 0.38, tamanho=17, cor=INK, negrito=True, fonte=FONT_HEAD)
    texto(s, "→ 13 casos Select AI medidos, limites explícitos e APEX real", 6.72, 2.62, 5.55, 0.58, tamanho=13, cor=GREEN, negrito=True)
    caixa(s, 0.82, 3.93, 5.55, 1.45, fill=AMBER_SOFT, line=AMBER, radius=True)
    rotulo(s, "Feedback da Sprint 1", 1.12, 4.24, 2.5, cor=AMBER)
    texto(s, "Mostrar protótipo funcional", 1.12, 4.71, 4.8, 0.38, tamanho=17, cor=INK, negrito=True, fonte=FONT_HEAD)
    texto(s, "→ WebApp público, quatro visões, ORDS live e reconciliação", 6.72, 4.43, 5.55, 0.58, tamanho=13, cor=GREEN, negrito=True)
    caixa(s, 1.45, 5.98, 10.45, 0.65, fill=INK, radius=True)
    texto(s, "Planejado na Sprint 1 → executado, publicado e auditável na Sprint 2", 1.72, 6.17, 9.92, 0.28, tamanho=13, cor=WHITE, negrito=True, alinhar=PP_ALIGN.CENTER)
    rodape(s, 18)

    # 19 — próximos passos
    s = prs.slides.add_slide(blank)
    fundo(s)
    titulo(s, "Próximos passos: a próxima evidência é humana", 19)
    caixa(s, 0.8, 2.22, 3.65, 3.82, fill=WHITE, line=AMBER, radius=True)
    rotulo(s, "Antes da banca", 1.1, 2.55, 2.3, cor=AMBER)
    lista(s, ["ensaiar a pergunta segura e o limite conhecido", "publicar vídeo e inserir o link", "fechar planilha e ZIP da entrega", "manter o Autonomous Database ativo"], 1.1, 3.08, 3.0, 2.2, tamanho=12.5, cor=INK, marcador="→", espacamento=8)
    caixa(s, 4.82, 2.22, 3.65, 3.82, fill=WHITE, line=GREEN, radius=True)
    rotulo(s, "Depois da entrega", 5.12, 2.55, 2.3)
    lista(s, ["entrevistas com gestores", "teste de compreensão dos indicadores", "teste de usabilidade da jornada", "piloto antes/depois"], 5.12, 3.08, 3.0, 2.2, tamanho=13, cor=INK, marcador="→", espacamento=8)
    caixa(s, 8.85, 2.22, 3.65, 3.82, fill=INK, radius=True)
    rotulo(s, "Pergunta de validação", 9.15, 2.55, 2.7, cor=MINT_DARK)
    texto(s, "O produto ajuda alguém\na chegar a uma hipótese\nmelhor, com menos erro\ne mais rastreabilidade?", 9.15, 3.15, 2.98, 2.05, tamanho=14.5, cor=WHITE, negrito=True, fonte=FONT_HEAD)
    texto(s, "Hoje, ainda não medimos.", 9.15, 5.32, 2.8, 0.35, tamanho=11, cor=MINT, negrito=True)
    rodape(s, 19)

    # 20 — equipe
    s = prs.slides.add_slide(blank)
    fundo(s, WHITE)
    titulo(s, "Uma equipe, uma cadeia de evidência", 20)
    nomes = [
        ("Carol Oliveira", "Bronze · preservação"),
        ("Lucas Lima", "Silver · Select AI"),
        ("Leandro Scutari", "Gold · indicadores"),
        ("Pedro Padovan", "Oracle · ORDS · frontend"),
        ("Leandro Lopes", "Oracle · ORDS · frontend"),
    ]
    for i, (nome, papel) in enumerate(nomes):
        y = 2.25 + i * 0.82
        texto(s, f"0{i + 1}", 0.85, y, 0.45, 0.28, tamanho=9, cor=GREEN, negrito=True)
        texto(s, nome, 1.55, y - 0.06, 3.0, 0.34, tamanho=18, cor=INK, negrito=True, fonte=FONT_HEAD)
        texto(s, papel, 4.75, y, 3.55, 0.3, tamanho=13, cor=MUTED)
        linha(s, 1.55, y + 0.5, 8.25, y + 0.5, LINE, 0.7)
    caixa(s, 8.85, 2.2, 3.55, 3.72, fill=INK, radius=True)
    texto(s, "Critério de pronto\ncompartilhado", 9.18, 2.58, 2.9, 1.02, tamanho=16.5, cor=WHITE, negrito=True, fonte=FONT_HEAD)
    lista(s, ["contrato", "reconciliação", "limite declarado", "evidência reproduzível"], 9.18, 3.78, 2.8, 1.65, tamanho=12, cor=MINT, marcador="✓", espacamento=6)
    rodape(s, 20)

    # 21 — conclusão e links
    s = prs.slides.add_slide(blank)
    fundo(s, INK)
    rotulo(s, "Conclusão · 21", 0.72, 0.55, 2.2, cor=MINT_DARK)
    texto(s, "Onde investigar primeiro?", 0.72, 1.12, 11.8, 0.7, tamanho=35, cor=WHITE, negrito=True, fonte=FONT_HEAD, alinhar=PP_ALIGN.CENTER)
    texto(s, "Do sinal territorial à hipótese: com fonte, amostra e limite à vista.", 0.72, 1.95, 11.8, 0.48, tamanho=15, cor=MINT_DARK, negrito=True, alinhar=PP_ALIGN.CENTER)
    texto(s, "O MedFlow não decide o que causou o sinal. Ele torna a próxima pergunta rastreável.", 1.15, 2.6, 11.0, 0.58, tamanho=15, cor=MINT, alinhar=PP_ALIGN.CENTER)
    linha(s, 0.72, 3.35, 12.3, 3.35, MINT_DARK, 1.2)
    links = [
        ("PRODUTO", "lucas-d-s-1.github.io/medflow", PRODUTO_URL, GREEN),
        ("REPOSITÓRIO", "github.com/Lucas-D-S-1/medflow", REPOSITORIO_URL, GREEN),
        ("VÍDEO", "YouTube · link pendente do time" if not VIDEO_URL else "Abrir demonstração no YouTube", VIDEO_URL, AMBER if not VIDEO_URL else GREEN),
    ]
    for i, (nome, endereco, destino, cor) in enumerate(links):
        x = 0.72 + i * 4.14
        caixa(s, x, 3.78, 3.78, 1.45, fill=MINT if destino else AMBER_SOFT, line=cor, radius=True)
        rotulo(s, nome, x + 0.25, 4.05, 2.8, cor=cor)
        texto(s, endereco, x + 0.25, 4.48, 3.28, 0.48, tamanho=10.5, cor=INK, negrito=True, alinhar=PP_ALIGN.CENTER)
        if destino:
            area_link(s, x, 3.78, 3.78, 1.45, destino)
    texto(s, "Carol · Leandro Lopes · Leandro Scutari · Lucas · Pedro", 0.75, 6.15, 11.85, 0.32, tamanho=10, cor=MINT, alinhar=PP_ALIGN.CENTER)
    texto(s, "FIAP × ORACLE · ÔMEGA URBAN TECH · 2026", 0.75, 6.68, 11.85, 0.28, tamanho=8.5, cor=MINT_DARK, negrito=True, alinhar=PP_ALIGN.CENTER)

    return prs


if __name__ == "__main__":
    apresentacao = criar()
    apresentacao.core_properties.title = "MedFlow — Sprint 2: evidências de construção"
    apresentacao.core_properties.subject = "Enterprise Challenge FIAP × Oracle 2026"
    apresentacao.core_properties.author = "Ômega Urban Tech"
    apresentacao.core_properties.comments = "Arco final de 21 slides baseado no produto e nas evidências vigentes."
    apresentacao.save(SAIDA)
    print(SAIDA)
